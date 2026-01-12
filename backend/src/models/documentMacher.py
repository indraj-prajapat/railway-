"""
COMPLETE FIXED VERSION: Enhanced DocumentMatcher with legend support for ALL file formats
Supports: PDF, DOCX, PPTX, XLSX, CSV, TXT with proper table caption/legend handling

Author: Enhanced version
Version: 3.0 - Complete legend support across all formats
"""
from __future__ import annotations

import io
import re
import logging
import typing as _t
from dataclasses import dataclass
from enum import Enum
from functools import lru_cache

logger = logging.getLogger(__name__)


# ---------- Types ------------------------------------------------------------
class MatchGranularity(str, Enum):
    LINE = "line"
    SENTENCE = "sentence"
    PARAGRAPH = "paragraph"
    CELL = "cell"
    TABLE_ROW = "table_row"
    SLIDE = "slide"  # For PPTX


@dataclass(slots=True)
class MatchSpan:
    text: str
    context: str
    granularity: MatchGranularity
    score: float
    location: dict[str, _t.Any]


@dataclass(slots=True)
class MatchResult:
    is_match: bool
    reason: str
    evidence: MatchSpan | None
    matches: list[MatchSpan]
    metadata: dict[str, _t.Any]


# ---------- Pure Helper Functions --------------------------------------------
def _normalise(text: str) -> str:
    """Lower-case + collapse whitespace for matching."""
    return re.sub(r"\s+", " ", text.lower())


@lru_cache(maxsize=2**14)
def _keyword_score_cached(norm_text: str, norm_query: str) -> float:
    """Cached keyword score – avoids re-counting identical strings."""
    if not norm_text or not norm_query:
        return 0.0
    return norm_text.count(norm_query) / (1 + len(norm_text.split()))


def _compute_dynamic_window(doc_length: int, granularity: MatchGranularity) -> int:
    """Compute context window size based on document length."""
    if granularity == MatchGranularity.LINE:
        if doc_length < 100:
            return 2
        elif doc_length < 500:
            return 3
        elif doc_length < 2000:
            return 5
        else:
            return 7
    elif granularity == MatchGranularity.SENTENCE:
        if doc_length < 50:
            return 1
        elif doc_length < 200:
            return 2
        else:
            return 3
    return 0


# ---------- Semantic Embedding (Optional) ------------------------------------
_EMB_CACHE: dict[str, _t.Any] = {}
_ST_MODEL: _t.Any = None
_ST_UTIL: _t.Any = None


def _get_st_model() -> tuple[_t.Any, _t.Any]:
    """Lazy load sentence-transformers model (optional)."""
    global _ST_MODEL, _ST_UTIL
    if _ST_MODEL is None:
        try:
            from sentence_transformers import SentenceTransformer, util
            _ST_MODEL = SentenceTransformer("all-MiniLM-L6-v2", device='cpu')
            _ST_UTIL = util
        except Exception as e:
            print(f"Semantic search disabled: {e}")
            return None, None
    return _ST_MODEL, _ST_UTIL


def _semantic_score(text: str, query: str) -> float:
    """Compute semantic similarity using cached embeddings."""
    if not text.strip():
        return 0.0
    
    try:
        model, util = _get_st_model()
        if model is None:
            return 0.0
        
        key = text
        if key not in _EMB_CACHE:
            _EMB_CACHE[key] = model.encode([text], convert_to_tensor=True, show_progress_bar=False)
        
        emb_q = model.encode([query], convert_to_tensor=True, show_progress_bar=False)
        similarity = float(util.cos_sim(emb_q, _EMB_CACHE[key])[0][0])
        return similarity
    except Exception:
        return 0.0


# ---------- File Type Detection ----------------------------------------------
def _detect_extension(data: bytes) -> str | None:
    """Detect file type from magic bytes."""
    if data.startswith(b"%PDF"):
        return ".pdf"
    if data.startswith(b"PK\x03\x04"):
        try:
            import zipfile
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                names = set(zf.namelist())
                if "word/document.xml" in names:
                    return ".docx"
                if "xl/workbook.xml" in names:
                    return ".xlsx"
                if "ppt/presentation.xml" in names:
                    return ".pptx"
        except Exception:
            pass
    return None


# ---------- Legend Parsing (UNIVERSAL) ---------------------------------------
def _parse_legend(text: str) -> dict[str, str]:
    """
    Extract legend mappings from any text.
    
    Patterns recognized:
    - "A = Bachelor of Science"
    - "A: Bachelor of Science"
    - "A - Bachelor of Science"
    - "1 = Footnote text"
    
    Returns dict of {key: description}
    """
    legend = {}
    
    # Pattern: Single char/digit = description
    patterns = [
        r'([A-Z0-9])\s*=\s*([^,;.\n]+)',  # A = Bachelor
        r'([A-Z0-9])\s*:\s*([^,;.\n]+)',  # A: Bachelor
        r'([A-Z0-9])\s+-\s+([^,;.\n]+)',  # A - Bachelor
        r'\(([0-9]+)\)\s*([^.\n]+)',      # (1) Footnote text
    ]
    
    for pattern in patterns:
        matches = re.finditer(pattern, text, re.IGNORECASE)
        for match in matches:
            key = match.group(1).strip().upper()
            value = match.group(2).strip()
            # Remove trailing punctuation
            value = re.sub(r'[,;.]$', '', value)
            if len(value) > 3:  # Avoid matching noise
                legend[key] = value
    
    return legend


def _expand_with_legend(text: str, legend: dict[str, str]) -> str:
    """
    Expand text by replacing legend keys with their definitions.
    
    Example:
        text = "A, C"
        legend = {"A": "Bachelor", "C": "Masters"}
        returns = "A (Bachelor), C (Masters)"
    """
    if not legend or not text:
        return text
    
    expanded_parts = []
    
    # Split by common separators
    parts = re.split(r'[,;\s]+', text)
    
    for part in parts:
        part = part.strip()
        if not part:
            continue
        
        # Check if this part is a legend key
        part_upper = part.upper()
        if part_upper in legend:
            expanded_parts.append(f"{part} ({legend[part_upper]})")
        else:
            expanded_parts.append(part)
    
    return ' '.join(expanded_parts)


def _extract_text_context(full_text: str, target_pos: int, window: int = 500) -> str:
    """Extract context around a position in text."""
    start = max(0, target_pos - window)
    end = min(len(full_text), target_pos + window)
    return full_text[start:end]


# ---------- PDF Processing ---------------------------------------------------
def _pdf_text_chunks(data: bytes) -> _t.Iterator[tuple[str, int]]:
    """Yield (page_text, page_number) – streaming, low memory."""
    try:
        import fitz
        doc = fitz.open(stream=data, filetype="pdf")
        for p, page in enumerate(doc, 1):
            yield page.get_text(), p
    except ImportError:
        import PyPDF2 as fpdf
        src = io.BytesIO(data)
        reader = fpdf.PdfReader(src)
        for i, page in enumerate(reader.pages, 1):
            yield (page.extract_text() or ""), i


def _extract_pdf_tables_with_context(data: bytes) -> list[dict[str, _t.Any]]:
    """Extract PDF tables with captions and surrounding context."""
    tables = []
    
    try:
        import pdfplumber
        import pandas as pd
        
        print("Extracting PDF tables with pdfplumber...")
        
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                page_tables = page.extract_tables()
                
                for table_idx, table in enumerate(page_tables):
                    if not table or len(table) < 2:
                        continue
                    
                    # Extract caption from page text (look for "Table N:")
                    caption_match = re.search(
                        rf'Table\s+{table_idx+1}[:\.]?\s*([^\n]+(?:\n[^\n]+)?)',
                        page_text,
                        re.IGNORECASE
                    )
                    caption = caption_match.group(0) if caption_match else ""
                    
                    # Parse legend from caption
                    legend = _parse_legend(caption)
                    
                    # Create DataFrame
                    headers = table[0]
                    data_rows = table[1:]
                    df = pd.DataFrame(data_rows, columns=headers)
                    df = df.dropna(how='all', axis=1).dropna(how='all', axis=0)
                    
                    if df.empty:
                        continue
                    
                    tables.append({
                        'df': df,
                        'page': page_num,
                        'table_index': table_idx,
                        'caption': caption,
                        'legend': legend,
                        'page_text': page_text,
                        'method': 'pdfplumber'
                    })
                    
                    print(f"PDF Table {table_idx+1} on page {page_num}: {df.shape}, legend={legend}")
        
        return tables
        
    except ImportError:
        print("pdfplumber not available, PDF table extraction disabled")
    except Exception as e:
        print(f"PDF table extraction failed: {e}")
    
    return []


def _highlight_pdf(data: bytes, query: str) -> bytes | None:
    """Add highlights to PDF for search query."""
    try:
        import fitz
        doc = fitz.open(stream=data, filetype="pdf")
        found = False
        
        for page in doc:
            for rect in page.search_for(query, quads=True):
                page.add_highlight_annot(rect).update()
                found = True
        
        if found:
            return doc.tobytes()
    except Exception:
        pass
    return None


# ---------- DOCX Processing --------------------------------------------------
def _docx_extract_all(data: bytes) -> dict[str, _t.Any]:
    """
    Extract text, tables, and styles from DOCX.
    Returns dict with 'text', 'tables', 'paragraphs'
    """
    from docx import Document
    
    doc = Document(io.BytesIO(data))
    
    # Extract paragraph text
    paragraphs = []
    full_text_parts = []
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append({
                'text': text,
                'style': para.style.name if para.style else None
            })
            full_text_parts.append(text)
    
    # Extract tables with context
    tables = []
    for table_idx, tbl in enumerate(doc.tables):
        # Get text before table (context)
        context = ""
        for para in paragraphs[-5:]:  # Last 5 paragraphs before table
            if 'table' in para['text'].lower():
                context = para['text']
                break
        
        # Parse legend from context
        legend = _parse_legend(context)
        
        # Extract table data
        table_data = []
        for row in tbl.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        
        if len(table_data) > 1:
            import pandas as pd
            headers = table_data[0]
            data_rows = table_data[1:]
            df = pd.DataFrame(data_rows, columns=headers)
            df = df.dropna(how='all', axis=1).dropna(how='all', axis=0)
            
            if not df.empty:
                tables.append({
                    'df': df,
                    'table_index': table_idx,
                    'context': context,
                    'legend': legend,
                    'raw_data': table_data
                })
                
                # Add table to full text
                full_text_parts.append(f"\n[Table {table_idx+1}]\n")
                full_text_parts.append("\t".join(headers))
                for row in data_rows:
                    full_text_parts.append("\t".join(row))
        
        print(f"DOCX Table {table_idx+1}: {df.shape if 'df' in locals() else 'empty'}, legend={legend}")
    
    return {
        'text': "\n".join(full_text_parts),
        'tables': tables,
        'paragraphs': paragraphs
    }


# ---------- PPTX Processing --------------------------------------------------
def _pptx_extract_all(data: bytes) -> dict[str, _t.Any]:
    """
    Extract text, tables, and notes from PPTX.
    Returns dict with 'slides', 'tables', 'text'
    """
    from pptx import Presentation
    
    prs = Presentation(io.BytesIO(data))
    slides = []
    tables = []
    full_text_parts = []
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_text = []
        slide_notes = ""
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
            
            # Extract tables
            if shape.has_table:
                tbl = shape.table
                
                # Look for legend in slide text
                slide_context = "\n".join(slide_text)
                legend = _parse_legend(slide_context)
                
                # Extract table data
                table_data = []
                for row in tbl.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                if len(table_data) > 1:
                    import pandas as pd
                    headers = table_data[0]
                    data_rows = table_data[1:]
                    df = pd.DataFrame(data_rows, columns=headers)
                    df = df.dropna(how='all', axis=1).dropna(how='all', axis=0)
                    
                    if not df.empty:
                        tables.append({
                            'df': df,
                            'slide_index': slide_idx,
                            'context': slide_context,
                            'legend': legend,
                            'raw_data': table_data
                        })
                        
                        print(f"PPTX Table on slide {slide_idx}: {df.shape}, legend={legend}")
        
        # Extract notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text:
                slide_notes = notes_frame.text
        
        slide_full_text = "\n".join(slide_text)
        slides.append({
            'index': slide_idx,
            'text': slide_full_text,
            'notes': slide_notes,
        })
        
        full_text_parts.append(f"--- Slide {slide_idx} ---")
        full_text_parts.append(slide_full_text)
        if slide_notes:
            full_text_parts.append(f"Notes: {slide_notes}")
    
    return {
        'text': "\n".join(full_text_parts),
        'slides': slides,
        'tables': tables
    }


# ---------- XLSX Processing --------------------------------------------------
def _xlsx_extract_all(data: bytes) -> dict[str, _t.Any]:
    """
    Extract all sheets from XLSX with legend detection.
    Returns dict with 'sheets' containing DataFrames and legends
    """
    import pandas as pd
    
    xls = pd.ExcelFile(io.BytesIO(data))
    sheets = {}
    
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name, dtype=str, keep_default_na=False)
        
        # Look for legend in first few rows
        legend = {}
        for i in range(min(5, len(df))):
            row_text = ' '.join(str(v) for v in df.iloc[i].values if str(v).strip())
            row_legend = _parse_legend(row_text)
            if row_legend:
                legend.update(row_legend)
        
        # Also check column names for legend
        header_text = ' '.join(str(col) for col in df.columns)
        legend.update(_parse_legend(header_text))
        
        sheets[sheet_name] = {
            'df': df,
            'legend': legend
        }
        
        print(f"XLSX Sheet '{sheet_name}': {df.shape}, legend={legend}")
    
    return {'sheets': sheets}


# ---------- CSV Processing ---------------------------------------------------
def _csv_extract(data: bytes) -> dict[str, _t.Any]:
    """Extract CSV data with legend detection."""
    import pandas as pd
    
    df = pd.read_csv(io.BytesIO(data), dtype=str, keep_default_na=False)
    
    # Look for legend in first few rows
    legend = {}
    for i in range(min(5, len(df))):
        row_text = ' '.join(str(v) for v in df.iloc[i].values if str(v).strip())
        row_legend = _parse_legend(row_text)
        if row_legend:
            legend.update(row_legend)
    
    # Check column names
    header_text = ' '.join(str(col) for col in df.columns)
    legend.update(_parse_legend(header_text))
    
    print(f"CSV: {df.shape}, legend={legend}")
    
    return {
        'df': df,
        'legend': legend
    }


# ---------- Plain Text Processing --------------------------------------------
def _plain_text(data: bytes) -> str:
    """Decode bytes to text."""
    try:
        return data.decode("utf-8")
    except Exception:
        return data.decode("latin-1", errors="ignore")


# ---------- Universal Table Search -------------------------------------------
def _search_table_universal(
    df: _t.Any,
    query: str,
    legend: dict[str, str],
    context: str = "",
    location_extra: dict = None
) -> list[MatchSpan]:
    """
    Universal table search with legend expansion.
    Works for any DataFrame from PDF, DOCX, PPTX, XLSX, CSV.
    """
    import pandas as pd
    
    # Tokenize query
    STOP_WORDS = {
        "the", "and", "or", "but", "in", "on", "at", "to", "for", "of", "with",
        "by", "a", "an", "is", "are", "was", "were", "be", "been", "have", "has",
        "had", "do", "does", "did", "will", "would", "could", "can", "should",
        "may", "might", "must", "shall", "my", "me", "i"
    }
    
    query_tokens = [
        t for t in re.findall(r'\w+', query.lower()) 
        if len(t) > 1 and t not in STOP_WORDS
    ]
    
    if not query_tokens:
        return []
    
    matches: list[MatchSpan] = []
    location_extra = location_extra or {}
    
    # Search each row
    for row_idx in df.index:
        row_data = df.iloc[row_idx]
        
        # Build searchable text with legend expansion
        original_parts = []
        expanded_parts = []
        
        for col_name, cell_value in row_data.items():
            if pd.isna(cell_value) or not str(cell_value).strip():
                continue
            
            cell_str = str(cell_value).strip()
            original_parts.append(f"{col_name}: {cell_str}")
            
            # Expand with legend
            expanded_cell = _expand_with_legend(cell_str, legend)
            if expanded_cell != cell_str:
                expanded_parts.append(f"{col_name}: {expanded_cell}")
        
        # Create searchable text
        original_text = " | ".join(original_parts)
        expanded_text = " | ".join(expanded_parts) if expanded_parts else ""
        searchable_text = f"{original_text} {expanded_text} {context}"
        searchable_lower = searchable_text.lower()
        
        # Check for keyword matches
        token_matches = sum(1 for token in query_tokens if token in searchable_lower)
        
        if token_matches == 0:
            continue
        
        # Compute score
        k_score = token_matches / len(query_tokens) if query_tokens else 0.0
        
        # Build context
        if expanded_parts:
            context_text = expanded_text
        else:
            context_text = original_text
        
        # Build location
        location = {
            'row_index': int(row_idx),
            'legend_used': bool(legend),
            **location_extra
        }
        
        matches.append(
            MatchSpan(
                text=context_text,
                context=context_text,
                granularity=MatchGranularity.TABLE_ROW,
                score=k_score,
                location=location
            )
        )
    
    # Sort by score
    matches.sort(key=lambda m: -m.score)
    return matches


# ---------- Core DocumentMatcher Class ---------------------------------------
class DocumentMatcher:
    """
    COMPLETE DocumentMatcher with legend support for ALL file formats.
    Supports: PDF, DOCX, PPTX, XLSX, CSV, TXT
    """

    def __init__(
        self,
        enable_semantic: bool = True,
        pdf_highlight: bool = False,
        pdf_table_extraction: bool = True,
        log_level: str = "INFO",
    ) -> None:
        self.enable_semantic = enable_semantic
        self.pdf_highlight = pdf_highlight
        self.pdf_table_extraction = pdf_table_extraction
        
        logging.basicConfig(
            level=getattr(logging, log_level.upper()),
            format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
        )
        
        print(
            f"DocumentMatcher (COMPLETE) initialized:\n"
            f"  - Semantic search: {enable_semantic}\n"
            f"  - PDF highlighting: {pdf_highlight}\n"
            f"  - Table extraction: {pdf_table_extraction}\n"
            f"  - Formats: PDF, DOCX, PPTX, XLSX, CSV, TXT"
        )

    def _best_spans(
        self, full_text: str, query: str, top_k: int = 8
    ) -> list[MatchSpan]:
        """Multi-granularity text matching with dynamic context windows."""
        norm_q = _normalise(query)
        if not norm_q:
            return []

        lines = full_text.splitlines()
        sentences = re.split(r"(?<=[.!?])\s+", full_text.strip())
        paras = [p.strip() for p in re.split(r"\n{2,}", full_text) if p.strip()]

        line_window = _compute_dynamic_window(len(lines), MatchGranularity.LINE)
        sent_window = _compute_dynamic_window(len(sentences), MatchGranularity.SENTENCE)

        spans: list[tuple[float, str, str, MatchGranularity, dict]] = []

        # Process lines
        for no, line in enumerate(lines, 1):
            norm_line = _normalise(line)
            if not norm_line:
                continue
            
            k_score = _keyword_score_cached(norm_line, norm_q)
            if k_score == 0 and not self.enable_semantic:
                continue
            
            s_score = _semantic_score(line, query) if self.enable_semantic else 0.0
            score = k_score + 0.3 * s_score
            
            if score <= 0:
                continue
            
            start_idx = max(0, no - line_window - 1)
            end_idx = min(len(lines), no + line_window)
            ctx = "\n".join(lines[start_idx:end_idx])
            
            spans.append((score, line, ctx, MatchGranularity.LINE, {"line_no": no}))

        # Process sentences
        for no, sent in enumerate(sentences, 1):
            norm_sent = _normalise(sent)
            if not norm_sent:
                continue
            
            k_score = _keyword_score_cached(norm_sent, norm_q)
            if k_score == 0 and not self.enable_semantic:
                continue
            
            s_score = _semantic_score(sent, query) if self.enable_semantic else 0.0
            score = k_score + 0.3 * s_score
            
            if score <= 0:
                continue
            
            start_idx = max(0, no - sent_window - 1)
            end_idx = min(len(sentences), no + sent_window)
            ctx = " ".join(sentences[start_idx:end_idx])
            
            spans.append((score, sent, ctx, MatchGranularity.SENTENCE, {"sentence_no": no}))

        # Process paragraphs
        for no, para in enumerate(paras, 1):
            norm_p = _normalise(para)
            k_score = _keyword_score_cached(norm_p, norm_q)
            if k_score == 0 and not self.enable_semantic:
                continue
            
            s_score = _semantic_score(para, query) if self.enable_semantic else 0.0
            score = k_score + 0.3 * s_score
            
            if score <= 0:
                continue
            
            spans.append((score, para, para, MatchGranularity.PARAGRAPH, {"paragraph_no": no}))

        spans.sort(reverse=True, key=lambda t: t[0])
        
        return [
            MatchSpan(text=t[1], context=t[2], granularity=t[3], score=t[0], location=t[4])
            for t in spans[:top_k]
        ]

    def match(
        self,
        file_input: bytes | io.IOBase,
        query: str,
        file_type: str | None = None,
    ) -> MatchResult:
        """
        Match query against document content.
        COMPLETE: Handles ALL file formats with legend support.
        """
        print(f"\n{'='*70}")
        print(f"DocumentMatcher Query: '{query}'")
        print(f"{'='*70}")
        
        # Normalize input
        if isinstance(file_input, io.IOBase):
            data = file_input.read()
        elif isinstance(file_input, bytes):
            data = file_input
        else:
            raise TypeError("file_input must be bytes or file-like object")

        # Detect file type
        detected = file_type or _detect_extension(data) or ".txt"
        detected = detected.lower()
        
        print(f"File type: {detected}, size: {len(data):,} bytes")

        q = query.strip()
        meta: dict[str, _t.Any] = {"file_type": detected, "query": q, "file_size": len(data)}
        matches: list[MatchSpan] = []
        evidence: MatchSpan | None = None
        is_match = False
        reason = "no_match"

        try:
            # ==================== PDF ====================
            if detected == ".pdf":
                print("\n=== Processing PDF ===")
                
                # Extract tables with legends
                if self.pdf_table_extraction:
                    tables = _extract_pdf_tables_with_context(data)
                    meta['tables_found'] = len(tables)
                    
                    for table_info in tables:
                        table_matches = _search_table_universal(
                            df=table_info['df'],
                            query=q,
                            legend=table_info['legend'],
                            context=table_info['caption'],
                            location_extra={
                                'page': table_info['page'],
                                'table_index': table_info['table_index'],
                                'source': 'pdf_table'
                            }
                        )
                        matches.extend(table_matches)
                        print(f"  Table {table_info['table_index']+1}: {len(table_matches)} matches")
                
                # Extract text content
                text_chunks = []
                for txt, page_num in _pdf_text_chunks(data):
                    text_chunks.append(txt)
                
                full_text = "\n".join(text_chunks)
                text_matches = self._best_spans(full_text, q, top_k=8)
                matches.extend(text_matches)
                print(f"  Text matches: {len(text_matches)}")
                
                # Determine reason
                if matches:
                    has_table = any(m.granularity == MatchGranularity.TABLE_ROW for m in matches)
                    reason = "table_match" if has_table else "keyword"
                
                # PDF highlighting
                if self.pdf_highlight and matches:
                    if highlighted := _highlight_pdf(data, q):
                        meta["pdf_highlighted_bytes"] = highlighted

            # ==================== DOCX ====================
            elif detected == ".docx":
                print("\n=== Processing DOCX ===")
                docx_data = _docx_extract_all(data)
                meta['tables_found'] = len(docx_data['tables'])
                
                # Search tables
                for table_info in docx_data['tables']:
                    table_matches = _search_table_universal(
                        df=table_info['df'],
                        query=q,
                        legend=table_info['legend'],
                        context=table_info['context'],
                        location_extra={
                            'table_index': table_info['table_index'],
                            'source': 'docx_table'
                        }
                    )
                    matches.extend(table_matches)
                    print(f"  Table {table_info['table_index']+1}: {len(table_matches)} matches")
                
                # Search text
                text_matches = self._best_spans(docx_data['text'], q, top_k=8)
                matches.extend(text_matches)
                print(f"  Text matches: {len(text_matches)}")
                
                if matches:
                    has_table = any(m.granularity == MatchGranularity.TABLE_ROW for m in matches)
                    reason = "table_match" if has_table else "keyword"

            # ==================== PPTX ====================
            elif detected == ".pptx":
                print("\n=== Processing PPTX ===")
                pptx_data = _pptx_extract_all(data)
                meta['slides_found'] = len(pptx_data['slides'])
                meta['tables_found'] = len(pptx_data['tables'])
                
                # Search tables
                for table_info in pptx_data['tables']:
                    table_matches = _search_table_universal(
                        df=table_info['df'],
                        query=q,
                        legend=table_info['legend'],
                        context=table_info['context'],
                        location_extra={
                            'slide_index': table_info['slide_index'],
                            'source': 'pptx_table'
                        }
                    )
                    matches.extend(table_matches)
                    print(f"  Slide {table_info['slide_index']} table: {len(table_matches)} matches")
                
                # Search text
                text_matches = self._best_spans(pptx_data['text'], q, top_k=8)
                matches.extend(text_matches)
                print(f"  Text matches: {len(text_matches)}")
                
                if matches:
                    has_table = any(m.granularity == MatchGranularity.TABLE_ROW for m in matches)
                    reason = "table_match" if has_table else "keyword"

            # ==================== XLSX ====================
            elif detected == ".xlsx":
                print("\n=== Processing XLSX ===")
                xlsx_data = _xlsx_extract_all(data)
                meta['sheets_found'] = len(xlsx_data['sheets'])
                
                for sheet_name, sheet_info in xlsx_data['sheets'].items():
                    sheet_matches = _search_table_universal(
                        df=sheet_info['df'],
                        query=q,
                        legend=sheet_info['legend'],
                        location_extra={
                            'sheet': sheet_name,
                            'source': 'xlsx_sheet'
                        }
                    )
                    matches.extend(sheet_matches)
                    print(f"  Sheet '{sheet_name}': {len(sheet_matches)} matches")
                
                if matches:
                    reason = "exact_cell_match"

            # ==================== CSV ====================
            elif detected == ".csv":
                print("\n=== Processing CSV ===")
                csv_data = _csv_extract(data)
                
                matches = _search_table_universal(
                    df=csv_data['df'],
                    query=q,
                    legend=csv_data['legend'],
                    location_extra={'source': 'csv'}
                )
                print(f"  Matches: {len(matches)}")
                
                if matches:
                    reason = "exact_cell_match"

            # ==================== Plain Text ====================
            else:
                print(f"\n=== Processing as Plain Text ({detected}) ===")
                full_text = _plain_text(data)
                matches = self._best_spans(full_text, q, top_k=8)
                print(f"  Matches: {len(matches)}")
                
                if matches:
                    reason = "keyword"

            # Finalize results
            if matches:
                is_match = True
                matches.sort(key=lambda m: -m.score)
                evidence = matches[0]
                print(f"\n✓ SUCCESS: {len(matches)} matches found")
                print(f"  Best score: {evidence.score:.4f}")
                print(f"  Granularity: {evidence.granularity.value}")
            else:
                print(f"\n✗ No matches found")

        except Exception as exc:
            print(f"\n✗ ERROR: {exc}")
            import traceback
            traceback.print_exc()
            reason = "error"
            meta["error"] = str(exc)

        print(f"{'='*70}\n")

        return MatchResult(
            is_match=is_match,
            reason=reason,
            evidence=evidence,
            matches=matches,
            metadata=meta,
        )


# ---------- Example Usage ----------------------------------------------------
if __name__ == "__main__":
    print("""
    ╔══════════════════════════════════════════════════════════════════════╗
    ║       COMPLETE DocumentMatcher with Universal Legend Support         ║
    ╚══════════════════════════════════════════════════════════════════════╝
    
    Supported formats with legend/table support:
    ✓ PDF    - Tables with captions and legends
    ✓ DOCX   - Tables with context and legends
    ✓ PPTX   - Slide tables with legends
    ✓ XLSX   - Multi-sheet with legend detection
    ✓ CSV    - With legend detection
    ✓ TXT    - Plain text search
    
    Features:
    • Automatic legend detection (A=Bachelor, C=Masters, etc.)
    • Cell value expansion ("A" → "A (Bachelor of Science)")
    • Context-aware table search
    • Multi-granularity text matching
    • Optional semantic similarity scoring
    • PDF highlighting support
    
    Installation:
      pip install pdfplumber pandas python-docx python-pptx openpyxl
      pip install sentence-transformers  # Optional, for semantic search
    
    Usage Example:
      matcher = DocumentMatcher(pdf_table_extraction=True)
      
      with open('document.pdf', 'rb') as f:
          result = matcher.match(f, 'masters course philosophy 2006')
      
      if result.is_match:
          print(f"Found {len(result.matches)} matches!")
          for match in result.matches[:3]:
              print(f"  {match.text[:80]}... (score: {match.score:.3f})")
    """)