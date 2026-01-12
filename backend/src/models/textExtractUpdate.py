"""
Enhanced text and structured data extraction from file bytes.
Extracts maximum information in a well-formatted way including tables.

Supports: PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, CSV, TXT
Features:
- Table extraction from PDFs
- Table preservation in DOCX/PPTX
- Structured data from spreadsheets
- Metadata inclusion
- Rich formatting preservation
"""
import io
import csv
import re
import logging
from typing import Dict, List, Any, Tuple, Optional

# Required imports
import PyPDF2
import docx
import pptx
import openpyxl
import xlrd
import filetype  # pip install filetype

# Optional imports for enhanced features
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    import tabula
    HAS_TABULA = True
except ImportError:
    HAS_TABULA = False

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False

logger = logging.getLogger(__name__)


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def _format_table_as_text(table_data: List[List[str]], title: str = "Table") -> str:
    """
    Format table data as aligned text with borders.
    
    Args:
        table_data: List of rows, each row is a list of cell values
        title: Title for the table
        
    Returns:
        Formatted table string
    """
    if not table_data:
        return ""
    
    # Calculate column widths
    col_widths = []
    for col_idx in range(len(table_data[0])):
        max_width = max(
            len(str(row[col_idx])) if col_idx < len(row) else 0 
            for row in table_data
        )
        col_widths.append(min(max_width, 50))  # Cap at 50 chars
    
    # Build formatted table
    lines = []
    lines.append(f"\n{'='*80}")
    lines.append(f"  {title}")
    lines.append('='*80)
    
    # Header row (if exists)
    if len(table_data) > 1:
        header = table_data[0]
        header_line = " | ".join(
            str(cell)[:width].ljust(width) 
            for cell, width in zip(header, col_widths)
        )
        lines.append(header_line)
        lines.append("-" * len(header_line))
    
    # Data rows
    for row in table_data[1:] if len(table_data) > 1 else table_data:
        row_line = " | ".join(
            str(cell)[:width].ljust(width) if idx < len(row) else "".ljust(width)
            for idx, (cell, width) in enumerate(zip(row, col_widths))
        )
        lines.append(row_line)
    
    lines.append('='*80)
    return "\n".join(lines)


def _dataframe_to_text(df, title: str = "Table") -> str:
    """Convert pandas DataFrame to formatted text."""
    if df.empty:
        return f"\n[{title}: Empty]\n"
    
    # Convert to list of lists
    table_data = [df.columns.tolist()] + df.values.tolist()
    return _format_table_as_text(table_data, title)


# ============================================================================
# PDF EXTRACTION WITH TABLES
# ============================================================================

def _extract_pdf_with_tables(file_bytes: bytes) -> str:
    """
    Extract text and tables from PDF with proper formatting.
    Uses pdfplumber if available, falls back to PyPDF2.
    """
    text_parts = []
    
    # Try pdfplumber first (best for tables)
    if HAS_PDFPLUMBER:
        try:
            logger.debug("Extracting PDF with pdfplumber")
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text_parts.append(f"\n{'='*80}\nPAGE {page_num}\n{'='*80}\n")
                    
                    # Extract tables first
                    tables = page.extract_tables()
                    if tables:
                        logger.debug(f"Found {len(tables)} tables on page {page_num}")
                        for table_idx, table in enumerate(tables, 1):
                            if table and len(table) > 0:
                                formatted_table = _format_table_as_text(
                                    table, 
                                    f"Table {table_idx} on Page {page_num}"
                                )
                                text_parts.append(formatted_table)
                    
                    # Extract regular text
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"\nText Content:\n{page_text}\n")
            
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning(f"pdfplumber extraction failed: {e}, falling back to PyPDF2")
    
    # Try tabula for tables if pdfplumber not available
    if HAS_TABULA and HAS_PANDAS:
        try:
            logger.debug("Extracting PDF tables with tabula")
            dfs = tabula.read_pdf(
                io.BytesIO(file_bytes),
                pages='all',
                multiple_tables=True,
                pandas_options={'dtype': str}
            )
            
            if dfs:
                logger.debug(f"Extracted {len(dfs)} tables with tabula")
                for idx, df in enumerate(dfs, 1):
                    if not df.empty:
                        formatted_table = _dataframe_to_text(df, f"Table {idx}")
                        text_parts.append(formatted_table)
        except Exception as e:
            logger.debug(f"tabula extraction failed: {e}")
    
    # Fallback: PyPDF2 for text only
    try:
        logger.debug("Extracting PDF text with PyPDF2")
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        for page_num, page in enumerate(reader.pages, 1):
            text_parts.append(f"\n{'='*80}\nPAGE {page_num}\n{'='*80}\n")
            page_text = page.extract_text() or ""
            text_parts.append(page_text)
    except Exception as e:
        logger.error(f"PyPDF2 extraction failed: {e}")
        text_parts.append(f"[Error extracting PDF: {e}]")
    
    return "\n".join(text_parts)


# ============================================================================
# DOCX EXTRACTION WITH TABLES
# ============================================================================

def _extract_docx_with_tables(file_bytes: bytes) -> str:
    """
    Extract text and tables from DOCX with proper formatting.
    Preserves document structure including tables.
    """
    text_parts = []
    
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        
        # Add document metadata
        core_props = doc.core_properties
        if core_props.title:
            text_parts.append(f"TITLE: {core_props.title}")
        if core_props.author:
            text_parts.append(f"AUTHOR: {core_props.author}")
        if core_props.subject:
            text_parts.append(f"SUBJECT: {core_props.subject}")
        
        if text_parts:
            text_parts.append("="*80)
        
        # Extract content in order (paragraphs and tables)
        table_count = 0
        for element in doc.element.body:
            # Check if it's a paragraph
            if element.tag.endswith('p'):
                para = None
                for p in doc.paragraphs:
                    if p._element == element:
                        para = p
                        break
                if para and para.text.strip():
                    text_parts.append(para.text)
            
            # Check if it's a table
            elif element.tag.endswith('tbl'):
                table_count += 1
                table = None
                for t in doc.tables:
                    if t._element == element:
                        table = t
                        break
                
                if table:
                    # Extract table data
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    if table_data:
                        formatted_table = _format_table_as_text(
                            table_data, 
                            f"Table {table_count}"
                        )
                        text_parts.append(formatted_table)
        
        logger.debug(f"Extracted {table_count} tables from DOCX")
        
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        text_parts.append(f"[Error extracting DOCX: {e}]")
    
    return "\n\n".join(text_parts)


# ============================================================================
# PPTX EXTRACTION WITH TABLES
# ============================================================================

def _extract_pptx_with_tables(file_bytes: bytes) -> str:
    """
    Extract text and tables from PPTX with proper formatting.
    Includes slide numbers and table formatting.
    """
    text_parts = []
    
    try:
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        
        # Add presentation metadata
        if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
            text_parts.append(f"PRESENTATION: {prs.core_properties.title}")
            text_parts.append("="*80)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n{'='*80}\nSLIDE {slide_num}\n{'='*80}\n")
            
            slide_text = []
            table_count = 0
            
            for shape in slide.shapes:
                # Extract text from text shapes
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract tables
                if shape.has_table:
                    table_count += 1
                    table = shape.table
                    
                    # Extract table data
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    if table_data:
                        formatted_table = _format_table_as_text(
                            table_data,
                            f"Table {table_count} on Slide {slide_num}"
                        )
                        slide_text.append(formatted_table)
            
            if slide_text:
                text_parts.append("\n".join(slide_text))
            
            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                text_parts.append(f"\nNotes:\n{slide.notes_slide.notes_text_frame.text}")
        
        logger.debug(f"Extracted from {len(prs.slides)} slides")
        
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        text_parts.append(f"[Error extracting PPTX: {e}]")
    
    return "\n\n".join(text_parts)


# ============================================================================
# XLSX EXTRACTION WITH FORMATTING
# ============================================================================

def _extract_xlsx_formatted(file_bytes: bytes) -> str:
    """
    Extract data from XLSX with proper table formatting for each sheet.
    Preserves structure and formatting.
    """
    text_parts = []
    
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        
        text_parts.append(f"EXCEL WORKBOOK: {len(wb.sheetnames)} sheets")
        text_parts.append("="*80)
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Get data range
            if ws.max_row == 0 or ws.max_column == 0:
                text_parts.append(f"\n[Sheet '{sheet_name}': Empty]\n")
                continue
            
            # Extract all data
            table_data = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                # Skip completely empty rows
                if any(cell.strip() for cell in row_data):
                    table_data.append(row_data)
            
            if table_data:
                formatted_table = _format_table_as_text(
                    table_data,
                    f"Sheet: {sheet_name}"
                )
                text_parts.append(formatted_table)
            
            logger.debug(f"Extracted {len(table_data)} rows from sheet '{sheet_name}'")
        
    except Exception as e:
        logger.error(f"XLSX extraction failed: {e}")
        text_parts.append(f"[Error extracting XLSX: {e}]")
    
    return "\n\n".join(text_parts)


# ============================================================================
# XLS EXTRACTION WITH FORMATTING
# ============================================================================

def _extract_xls_formatted(file_bytes: bytes) -> str:
    """
    Extract data from XLS (old Excel) with proper table formatting.
    """
    text_parts = []
    
    try:
        wb = xlrd.open_workbook(file_contents=file_bytes)
        
        text_parts.append(f"EXCEL WORKBOOK (Legacy): {wb.nsheets} sheets")
        text_parts.append("="*80)
        
        for sheet in wb.sheets():
            if sheet.nrows == 0 or sheet.ncols == 0:
                text_parts.append(f"\n[Sheet '{sheet.name}': Empty]\n")
                continue
            
            # Extract all data
            table_data = []
            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)
                row_data = [str(cell) if cell else "" for cell in row]
                # Skip completely empty rows
                if any(cell.strip() for cell in row_data):
                    table_data.append(row_data)
            
            if table_data:
                formatted_table = _format_table_as_text(
                    table_data,
                    f"Sheet: {sheet.name}"
                )
                text_parts.append(formatted_table)
            
            logger.debug(f"Extracted {len(table_data)} rows from sheet '{sheet.name}'")
        
    except Exception as e:
        logger.error(f"XLS extraction failed: {e}")
        text_parts.append(f"[Error extracting XLS: {e}]")
    
    return "\n\n".join(text_parts)


# ============================================================================
# CSV EXTRACTION WITH FORMATTING
# ============================================================================

def _extract_csv_formatted(file_bytes: bytes) -> str:
    """
    Extract data from CSV with proper table formatting.
    """
    text_parts = []
    
    try:
        # Try different encodings
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                text_content = file_bytes.decode(encoding, errors='ignore')
                break
            except:
                continue
        else:
            text_content = file_bytes.decode('utf-8', errors='ignore')
        
        csvfile = io.StringIO(text_content)
        reader = csv.reader(csvfile)
        
        table_data = []
        for row in reader:
            if any(cell.strip() for cell in row):  # Skip empty rows
                table_data.append(row)
        
        if table_data:
            formatted_table = _format_table_as_text(table_data, "CSV Data")
            text_parts.append(formatted_table)
        
        logger.debug(f"Extracted {len(table_data)} rows from CSV")
        
    except Exception as e:
        logger.error(f"CSV extraction failed: {e}")
        text_parts.append(f"[Error extracting CSV: {e}]")
    
    return "\n".join(text_parts)


# ============================================================================
# MAIN EXTRACTION FUNCTION
# ============================================================================

def extract_text_from_bytes(
    file_bytes: bytes, 
    mime: Optional[str] = None, 
    filename: Optional[str] = None
) -> str:
    """
    Extract text and structured data from raw file bytes.
    
    Features:
    - Extracts tables from PDFs with proper formatting
    - Preserves tables in DOCX and PPTX
    - Formats spreadsheet data as tables
    - Includes metadata when available
    - Handles multiple file formats
    
    Supports: PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, CSV, TXT
    
    Args:
        file_bytes: Raw file content as bytes
        mime: MIME type (optional, helps with detection)
        filename: Original filename (optional, appended to output)
    
    Returns:
        Extracted text with formatted tables and structure
    """
    text = ""
    
    try:
        # Detect file type
        kind = filetype.guess(file_bytes)
        ext = kind.extension if kind else None
        
        logger.debug(f"Detected file type: {ext}, MIME: {mime}")
        
        # --- PDF (ENHANCED with tables) ---
        if mime == "application/pdf" or ext == "pdf":
            text = _extract_pdf_with_tables(file_bytes)
        
        # --- DOCX (ENHANCED with tables) ---
        elif ext == "docx" or mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text = _extract_docx_with_tables(file_bytes)
        
        # --- DOC (old Word) ---
        elif ext == "doc" or mime == "application/msword":
            try:
                import textract
                text = textract.process(io.BytesIO(file_bytes)).decode("utf-8", errors="ignore")
            except ImportError:
                logger.warning("textract not available for .doc files")
                text = "[.doc format requires 'textract' package: pip install textract]"
            except Exception as e:
                text = f"[Error extracting .doc file: {e}]"
        
        # --- PPTX (ENHANCED with tables) ---
        elif ext == "pptx" or mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = _extract_pptx_with_tables(file_bytes)
        
        # --- PPT (old PowerPoint) ---
        elif ext == "ppt" or mime == "application/vnd.ms-powerpoint":
            try:
                import textract
                text = textract.process(io.BytesIO(file_bytes)).decode("utf-8", errors="ignore")
            except ImportError:
                logger.warning("textract not available for .ppt files")
                text = "[.ppt format requires 'textract' package: pip install textract]"
            except Exception as e:
                text = f"[Error extracting .ppt file: {e}]"
        
        # --- XLSX (ENHANCED with formatting) ---
        elif ext == "xlsx" or mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            text = _extract_xlsx_formatted(file_bytes)
        
        # --- XLS (ENHANCED with formatting) ---
        elif ext == "xls" or mime == "application/vnd.ms-excel":
            text = _extract_xls_formatted(file_bytes)
        
        # --- CSV (ENHANCED with formatting) ---
        elif mime == "text/csv" or (filename and filename.endswith('.csv')):
            text = _extract_csv_formatted(file_bytes)
        
        # --- TXT and other text formats ---
        elif ext in ["txt", "log", "md", "json", "xml"] or (mime and mime.startswith("text/")):
            try:
                text = file_bytes.decode("utf-8", errors="ignore")
            except:
                text = file_bytes.decode("latin-1", errors="ignore")
        
        else:
            text = f"[Unsupported file type: {ext or mime or 'unknown'}]"
            logger.warning(f"Unsupported file type: {ext or mime}")
        
        # Append filename if provided
        if filename:
            text = f"FILENAME: {filename}\n{'='*80}\n\n{text}"
        
    except Exception as e:
        logger.error(f"Error extracting text: {e}", exc_info=True)
        text = f"[Error extracting text: {str(e)}]"
    
    return text.strip()


# ============================================================================
# USAGE EXAMPLES
# ============================================================================

if __name__ == "__main__":
    # Configure logging
    logging.basicConfig(level=logging.INFO)
    
    print("Enhanced Text Extraction with Tables")
    print("="*80)
    print("\nFeatures:")
    print("✓ PDF table extraction with pdfplumber/tabula")
    print("✓ DOCX table preservation")
    print("✓ PPTX table extraction")
    print("✓ Formatted XLSX/XLS/CSV output")
    print("✓ Metadata inclusion")
    print("✓ Rich structure preservation")
    
    print("\n" + "="*80)
    print("Usage Example:")
    print("="*80)
    print("""
# Read file
with open("document.pdf", "rb") as f:
    file_bytes = f.read()

# Extract text and tables
text = extract_text_from_bytes(
    file_bytes, 
    mime="application/pdf",
    filename="document.pdf"
)

print(text)

# Output will include:
# - Page numbers
# - Formatted tables with borders
# - Regular text content
# - Metadata (if available)
""")
    
    print("\n" + "="*80)
    print("Installation:")
    print("="*80)
    print("""
# Required
pip install PyPDF2 python-docx python-pptx openpyxl xlrd filetype

# Recommended for PDF tables
pip install pdfplumber pandas

# Alternative for PDF tables (requires Java)
pip install tabula-py pandas

# Optional for old Office formats (.doc, .ppt)
pip install textract
""")
    
    print("\n" + "="*80)
    print("Comparison:")
    print("="*80)
    print("""
BEFORE (Original):
------------------
Quarter Revenue Expenses Q1 100K 80K Q2 120K 90K

AFTER (Enhanced):
-----------------
================================================================================
  Table 1 on Page 1
================================================================================
Quarter    | Revenue    | Expenses
------------------------------------------
Q1         | 100K       | 80K
Q2         | 120K       | 90K
================================================================================
""")