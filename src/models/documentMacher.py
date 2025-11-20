"""
High-performance drop-in replacement for the original DocumentMatcher.
Author:  <you>
"""
from __future__ import annotations

import io
import os
import re
import json
import typing as _t
from dataclasses import dataclass
from enum import Enum
from functools import lru_cache

# ---------- types ------------------------------------------------------------
class MatchGranularity(str, Enum):
    LINE = "line"
    SENTENCE = "sentence"
    PARAGRAPH = "paragraph"
    CELL = "cell"

@dataclass(slots=True)
class MatchSpan:
    text: str
    context: str
    granularity: MatchGranularity
    score: float
    location: dict[str, _t.Any]

@dataclass(slots=True)
class MatchResult:
    # is_match: bool
    reason: str
    evidence: MatchSpan | None
    # matches: list[MatchSpan]
    # metadata: dict[str, _t.Any]

# ---------- small pure helpers ----------------------------------------------
def _normalise(text: str) -> str:
    """Lower-case + collapse whitespace for matching."""
    return re.sub(r"\s+", " ", text.lower())

@lru_cache(maxsize=2**14)
def _keyword_score_cached(norm_text: str, norm_query: str) -> float:
    """Cached keyword score – avoids re-counting identical strings."""
    return norm_text.count(norm_query) / (1 + len(norm_text.split()))

# ---------- semantic cache ---------------------------------------------------
# We cache the *embedding* of a text, not the similarity, so we can reuse
# the same embedding for every query.
_EMB_CACHE: dict[str, _t.Any] = {}  # text -> vector
_ST_MODEL: _t.Any = None
_ST_UTIL: _t.Any = None

def _get_st_model() -> _t.Any:
    global _ST_MODEL, _ST_UTIL
    if _ST_MODEL is None:
        from sentence_transformers import SentenceTransformer, util
        _ST_MODEL = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")
        _ST_UTIL = util
    return _ST_MODEL, _ST_UTIL

def _semantic_score(text: str, query: str) -> float:
    if not text.strip():
        return 0.0
    model, util = _get_st_model()
    key = text
    if key not in _EMB_CACHE:
        _EMB_CACHE[key] = model.encode([text], convert_to_tensor=True)
    emb_q = model.encode([query], convert_to_tensor=True)
    return float(util.cos_sim(emb_q, _EMB_CACHE[key])[0][0])

# ---------- Aho-Corasick keyword scanner -------------------------------------
try:
    import ahocorasick  # type: ignore
    _AC_AVAILABLE = True
except ImportError:
    _AC_AVAILABLE = False

def _build_ac(keys: list[str]) -> ahocorasick.Automaton | None:
    if not _AC_AVAILABLE or not keys:
        return None
    A = ahocorasick.Automaton()
    for k in keys:
        A.add_word(k.lower(), k.lower())
    A.make_automaton()
    return A

# ---------- file-type detection ---------------------------------------------
def _detect_extension(data: bytes) -> str | None:
    if data.startswith(b"%PDF"):
        return ".pdf"
    if data.startswith(b"PK\x03\x04"):
        try:
            import zipfile
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                names = set(zf.namelist())
                if "word/document.xml" in names:
                    return ".docx"
                if "xl/workbook.xml" in names:
                    return ".xlsx"
                if "ppt/presentation.xml" in names:
                    return ".pptx"
        except Exception:
            pass
    return None

# ---------- lazy text extractors --------------------------------------------
def _pdf_text_chunks(data: bytes) -> _t.Iterator[tuple[str, int]]:
    """Yield (page_text, page_number) – streaming, low memory."""
    try:
        import fitz  # type: ignore
    except ImportError:
        import PyPDF2 as fpdf

        src = io.BytesIO(data)
        reader = fpdf.PdfReader(src)
        for i, page in enumerate(reader.pages, 1):
            yield (page.extract_text() or ""), i
        return

    doc = fitz.open(stream=data, filetype="pdf")
    for p, page in enumerate(doc, 1):
        yield page.get_text(), p

def _docx_text(data: bytes) -> str:
    from docx import Document  # type: ignore
    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for tbl in doc.tables:
        for row in tbl.rows:
            parts.append("\t".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)

def _pptx_text(data: bytes) -> str:
    from pptx import Presentation  # type: ignore
    prs = Presentation(io.BytesIO(data))
    out: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        out.append(f"--- Slide {idx} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                out.append(shape.text)
    return "\n".join(out)

def _xlsx_tables(data: bytes) -> dict[str, _t.Any]:
    import pandas as pd
    xls = pd.ExcelFile(io.BytesIO(data))
    return {name: xls.parse(name, dtype=str, keep_default_na=False) for name in xls.sheet_names}

def _csv_table(data: bytes) -> _t.Any:
    import pandas as pd
    return pd.read_csv(io.BytesIO(data), dtype=str, keep_default_na=False)

def _plain_text(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except Exception:
        return data.decode("latin-1", errors="ignore")
from typing import Any, List
# ---------- core high-speed matcher -----------------------------------------
class DocumentMatcher:
    def __init__(
        self,
        enable_semantic: bool = True,
        pdf_highlight: bool = False,
    ) -> None:
        self.enable_semantic = enable_semantic
        self.pdf_highlight = pdf_highlight
        # We keep the model loading lazy – first call to _semantic_score
        # will import sentence-transformers if enable_semantic is True.

    # ---- internal helpers ----------------------------------------------------
    def _best_spans(
        self, full_text: str, query: str, top_k: int = 8
    ) -> list[MatchSpan]:
        """Multi-granularity matching – single pass, cached embeddings."""
        norm_q = _normalise(query)
        if not norm_q:
            return []

        # Build candidate list once
        lines = full_text.splitlines()
        sentences = re.split(r"(?<=[.!?])\s+", full_text.strip())
        paras = [p.strip() for p in re.split(r"\n{2,}", full_text) if p.strip()]

        spans: list[tuple[float, str, str, MatchGranularity, dict]] = []

        # Line level
        for no, line in enumerate(lines, 1):
            if not (norm_line := _normalise(line)):
                continue
            k_score = _keyword_score_cached(norm_line, norm_q)
            if k_score == 0 and not self.enable_semantic:
                continue
            s_score = _semantic_score(line, query) if self.enable_semantic else 0.0
            score = k_score + 0.3 * s_score
            if score <= 0:
                continue
            ctx = "\n".join(lines[max(0, no - 3) : no + 2])
            spans.append((score, line, ctx, MatchGranularity.LINE, {"line_no": no}))

        # Sentence level
        for no, sent in enumerate(sentences, 1):
            if not (norm_sent := _normalise(sent)):
                continue
            k_score = _keyword_score_cached(norm_sent, norm_q)
            if k_score == 0 and not self.enable_semantic:
                continue
            s_score = _semantic_score(sent, query) if self.enable_semantic else 0.0
            score = k_score + 0.3 * s_score
            if score <= 0:
                continue
            ctx = " ".join(sentences[max(0, no - 3) : no + 2])
            spans.append((score, sent, ctx, MatchGranularity.SENTENCE, {"sentence_no": no}))

        # Paragraph level
        for no, para in enumerate(paras, 1):
            norm_p = _normalise(para)
            k_score = _keyword_score_cached(norm_p, norm_q)
            if k_score == 0 and not self.enable_semantic:
                continue
            s_score = _semantic_score(para, query) if self.enable_semantic else 0.0
            score = k_score + 0.3 * s_score
            if score <= 0:
                continue
            spans.append((score, para, para, MatchGranularity.PARAGRAPH, {"paragraph_no": no}))

        spans.sort(reverse=True, key=lambda t: t[0])
        return [
            MatchSpan(text=t[1], context=t[2], granularity=t[3], score=t[0], location=t[4])
            for t in spans[:top_k]
        ]

    def _search_table(self, df: _t.Any, query: str, sheet: str | None = None) -> list[MatchSpan]:
        """Vectorised cell search – falls back to iterrows only if necessary."""
        import pandas as pd
        ql = query.lower()
        STOP_WORDS = {
            "the", "and", "or", "but", "in", "on", "at", "to", "for", "of", "with", "by", "a", "an", "is", "are", "was", "were", "be", "been", "have", "has", "had", "do", "does", "did", "will", "would", "could", "can", "should", "may", "might", "must", "shall"
        }
        def _normalise(txt: str) -> List[str]:
            """Lower-case, keep only alphanumerics, drop 1-letter tokens & stop-words."""
            txt = re.sub(r"[^a-z0-9]+", " ", txt.lower())
            tokens = {t for t in txt.split() if len(t) > 1} - STOP_WORDS
            return list(tokens)
        ql = _normalise(ql)
        print('dataframe',df)
        print('query lower',ql)
        print('sheet name',sheet)
        # Fast vectorised mask
        try:
            word_masks = {
                w: df.map(lambda cell, word=w: word in str(cell).lower())
                for w in ql
            }
            row_mask = pd.concat(word_masks.values(), axis=1).groupby(level=0).any()
            print('mask',row_mask)
            hit_cells = row_mask.stack()[row_mask.stack()]   # Series whose index is (row, col_name)
            hit_rows  = hit_cells.index.tolist() 
        except Exception as e:
            print("!!!!!! vectorised search exception !!!!!!", e)
            # Fallback to slower per-row scan
           

        print('hit rows',hit_rows)
        col_name_to_pos = {name: pos for pos, name in enumerate(df.columns)}
        if not hit_rows:
            return []

        matches: list[MatchSpan] = []
        import traceback,sys
        for (r, col_name) in hit_rows:
            try:
                c = col_name_to_pos[col_name] 
                print('r, c ',r , c)

                cell = str(df.iat[r, c])
                print('cell',cell)
                score = 0.0
                for q in ql:
                    score += cell.lower().count(q)
                score = score/(1 + len(cell.split()))
                print('score',score)
                ctx = " | ".join(f"{k}={v}" for k, v in list(df.iloc[r, :].items())[:5])
                print('context',ctx)
                matches.append(
                    MatchSpan(
                        text=cell,
                        context=f"Row {r} • {ctx}",
                        granularity=MatchGranularity.CELL,
                        score=score,
                        location={"row_index": int(r), "col_index": int(c), "col_name": str(c), **({"sheet": sheet} if sheet else {})},
                    )
                )
            except Exception as inner:  
                print("!!!!!! inner loop exception !!!!!!", file=sys.stderr)
                traceback.print_exception(inner)
        matches.sort(key=lambda m: (-m.score, m.location.get("row_index", 0)))
        print('matches',matches)
        return matches

    def _highlight_pdf(self, data: bytes, query: str) -> bytes | None:
        try:
            import fitz  # type: ignore
        except ImportError:
            return None
        doc = fitz.open(stream=data, filetype="pdf")
        found = False
        for page in doc:
            for rect in page.search_for(query, quads=True):
                page.add_highlight_annot(rect).update()
                found = True
        return doc.tobytes() if found else None

    # ---- public API ----------------------------------------------------------
    def match(
        self,
        file_input: bytes | io.IOBase,
        query: str,
        file_type: str | None = None,
    ) -> MatchResult:
        # normalise input
        if isinstance(file_input, io.IOBase):
            data = file_input.read()
            if isinstance(file_input, io.BytesIO):
                file_input.seek(0)
        elif isinstance(file_input, bytes):
            data = file_input
        else:
            raise TypeError("file_input must be bytes or file-like object")
 
        detected = file_type or _detect_extension(data) or ".txt"
        detected = detected.lower()
        print('detected',detected)
        q = query.strip()
        meta: dict[str, _t.Any] = {"file_type": detected, "query": q}
        matches: list[MatchSpan] = []
        evidence: MatchSpan | None = None
        is_match = False
        reason = "no_match"

        try:
            if detected == ".pdf":
                try:
                    text_chunks: list[str] = []
                    for txt, _ in _pdf_text_chunks(data):
                        text_chunks.append(txt)
                    full_text = "\n".join(text_chunks)
                    matches = self._best_spans(full_text, q, top_k=8)
                    reason = "keyword" if matches else "no_match"
                    if self.pdf_highlight and matches:
                        if highlighted := self._highlight_pdf(data, q):
                            meta["pdf_highlighted_bytes"] = highlighted
                except Exception as e:
                    print("PDF processing error:", e)
                    reason = "error"
                    meta["error"] = str(e)

            elif detected == ".docx":
                full_text = _docx_text(data)
                matches = self._best_spans(full_text, q, top_k=8)
                reason = "keyword" if matches else "no_match"

            elif detected == ".pptx":
                full_text = _pptx_text(data)
                matches = self._best_spans(full_text, q, top_k=8)
                reason = "keyword" if matches else "no_match"

            elif detected == ".xlsx":
                sheets = _xlsx_tables(data)
                for name, df in sheets.items():
                 
                    matches.extend(self._search_table(df, q, sheet=name))
                reason = "exact_cell_match" if matches else "no_match"
                meta.update({"sheet_count": len(sheets), "hit_count": len(matches)})

            elif detected == ".csv":
                df = _csv_table(data)
                matches = self._search_table(df, q)
                reason = "exact_cell_match" if matches else "no_match"
                meta["hit_count"] = len(matches)

            else:  # treat as plain text
                full_text = _plain_text(data)
                matches = self._best_spans(full_text, q, top_k=8)
                reason = "keyword" if matches else "no_match"

            if matches:
                is_match = True
                evidence = matches[0]

        except Exception as exc:
            reason = "error"
            meta["error"] = str(exc)

        return MatchResult(
            # is_match=is_match,
            reason=reason,
            evidence=evidence,
            # matches=matches,
            # metadata=meta,
        )